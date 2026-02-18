from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt


def main() -> None:
    img_paths = [
        r"C:\Users\Hasee\.cursor\projects\d-Daew-OCR-Tool-2\assets\c__Users_Hasee_AppData_Roaming_Cursor_User_workspaceStorage_858280dba88950aa495773d4845d8d44_images_image-3249cd5c-4a73-4ab5-8525-6c8f865aeff4.png",
        r"C:\Users\Hasee\.cursor\projects\d-Daew-OCR-Tool-2\assets\c__Users_Hasee_AppData_Roaming_Cursor_User_workspaceStorage_858280dba88950aa495773d4845d8d44_images_image-905c4100-fd28-468f-9b6b-bc0d58c45c16.png",
        r"C:\Users\Hasee\.cursor\projects\d-Daew-OCR-Tool-2\assets\c__Users_Hasee_AppData_Roaming_Cursor_User_workspaceStorage_858280dba88950aa495773d4845d8d44_images_image-acb42a4e-66d0-465d-ae1c-55c3feea389b.png",
    ]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.8), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.text = "Industrial OCR Measurement Tool"
    for p in title_tf.paragraphs:
        p.font.size = Pt(40)
        p.font.bold = True

    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.8), Inches(0.6))
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.text = "Progress summary and next steps"
    for p in subtitle_tf.paragraphs:
        p.font.size = Pt(20)

    # Slide 2: Project overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    hdr = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.6))
    hdr_tf = hdr.text_frame
    hdr_tf.text = "Project Overview"
    hdr_tf.paragraphs[0].font.size = Pt(30)
    hdr_tf.paragraphs[0].font.bold = True

    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12), Inches(5.6))
    tf = box.text_frame
    bullets = [
        "Goal: extract equipment dimensions from images and compute surface area.",
        "Focus: industrial containers, buckets, scoops, cylinders, and frustums.",
        "Output: labeled images, detected dimensions, and calculated areas.",
    ]

    tf.text = bullets[0]
    tf.paragraphs[0].font.size = Pt(18)
    for text in bullets[1:]:
        p = tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(18)

    # Slide 3: What we built
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    hdr = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.6))
    hdr_tf = hdr.text_frame
    hdr_tf.text = "What We Built"
    hdr_tf.paragraphs[0].font.size = Pt(30)
    hdr_tf.paragraphs[0].font.bold = True

    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12), Inches(5.6))
    tf = box.text_frame
    bullets = [
        "Local web app to upload images and run OCR.",
        "Visualization output with bounding boxes and labels.",
        "Surface-area calculator for common equipment shapes.",
        "Manual calculator page for exact measurements.",
    ]
    tf.text = bullets[0]
    tf.paragraphs[0].font.size = Pt(18)
    for text in bullets[1:]:
        p = tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(18)

    # Slide 4: Sample results (with image)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    hdr = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.6))
    hdr_tf = hdr.text_frame
    hdr_tf.text = "Sample Results"
    hdr_tf.paragraphs[0].font.size = Pt(30)
    hdr_tf.paragraphs[0].font.bold = True

    img = Path(img_paths[0])
    if img.exists():
        slide.shapes.add_picture(str(img), Inches(0.8), Inches(1.2), width=Inches(11.8))

    # Slide 5: Next steps
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    hdr = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.6))
    hdr_tf = hdr.text_frame
    hdr_tf.text = "Next Steps"
    hdr_tf.paragraphs[0].font.size = Pt(30)
    hdr_tf.paragraphs[0].font.bold = True

    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12), Inches(5.6))
    tf = box.text_frame
    bullets = [
        "I will host the app online so more people can use it easily.",
        "Improve OCR accuracy on new dataset images.",
        "Add more equipment types and measurement rules.",
        "Collect feedback and iterate on the UI.",
    ]
    tf.text = bullets[0]
    tf.paragraphs[0].font.size = Pt(18)
    for text in bullets[1:]:
        p = tf.add_paragraph()
        p.text = text
        p.level = 0
        p.font.size = Pt(18)

    out_path = Path(r"d:\Daew\OCR Tool 2\OCR_Tool_Presentation.pptx")
    prs.save(out_path)


if __name__ == "__main__":
    main()
